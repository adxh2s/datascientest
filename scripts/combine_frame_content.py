import json
import os
from pptx import Presentation
from pptx.util import Inches

DATA_DIR = "data"
OUTPUT_PPTX = "rakuten_presentation.pptx"

def load_json(path):
    with open(path, 'r', encoding='utf-8') as f:
        return json.load(f)

def apply_layout(slide, layout_data):
    slide_layout = slide.slide_layout
    # (python-pptx: changer le layout via slide_layouts)
    if 'background_color' in layout_data:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = layout_data['background_color']
    if 'title' in layout_data:
        title = slide.shapes.title
        title.text = layout_data['title']['text']
        title.text_frame.paragraphs[0].font.color.rgb = layout_data['title']['color']
        title.text_frame.paragraphs[0].font.size = layout_data['title']['size']

def apply_content(slide, content_data):
    if content_data['type'] == 'text':
        tf = slide.shapes.placeholders[1].text_frame
        tf.text = content_data['content']
        tf.paragraphs[0].font.color.rgb = content_data['style']['color']
        tf.paragraphs[0].font.size = content_data['style']['size']
    elif content_data['type'] == 'table':
        data = content_data['data']
        rows = len(data)
        cols = len(data[0])
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(2)
        table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
        for i, row in enumerate(data):
            for j, cell in enumerate(row):
                table.cell(i, j).text = str(cell)
    elif content_data['type'] == 'image':
        img_path = os.path.join(DATA_DIR, content_data['path'])
        slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(content_data['width']/72), height=Inches(content_data['height']/72))

def generate_pptx():
    prs = Presentation()
    # Liste des fichiers JSON pour chaque diapositive
    slide_files = sorted([f for f in os.listdir(DATA_DIR) if f.startswith('rakuten_presentation_slide') and f.endswith('.json')])
    slides = {}
    for f in slide_files:
        parts = f.split('_')
        num = parts[2][5:]  # extraire le numéro (slide01 -> 01)
        part = parts[3].split('.')[0]
        if num not in slides:
            slides[num] = {}
        slides[num][part] = os.path.join(DATA_DIR, f)
    # Parcourir chaque diapositive
    for num in sorted(slides.keys()):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # layout "Title and Content"
        # Appliquer le cadre
        if 'cadre' in slides[num]:
            cadre = load_json(slides[num]['cadre'])
            apply_layout(slide, cadre)
        # Appliquer le contenu
        if 'contenu' in slides[num]:
            contenu = load_json(slides[num]['contenu'])
            apply_content(slide, contenu)
    prs.save(OUTPUT_PPTX)
    print(f"Présentation générée : {OUTPUT_PPTX}")

if __name__ == "__main__":
    generate_pptx()
