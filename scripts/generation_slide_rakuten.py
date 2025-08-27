from pptx import Presentation
from pptx.util import Inches

# Création d'une présentation vide
prs = Presentation()

# 1. Diapositive de titre
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Rakuten : Synthèse et stratégie"
subtitle.text = "Focus France & panorama concurrentiel"

# 2. Informations marquantes sur Rakuten
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "1. Informations marquantes sur Rakuten"
content.text = (
    "- Chiffre d’affaires consolidé : 2 279 milliards de yens (environ 13,4 milliards d’euros), +10 %\n"
    "- Croissance des trois divisions : e-commerce, services financiers, téléphonie mobile\n"
    "- Résultat opérationnel positif pour la première fois depuis 2020\n"
    "- Perte nette réduite de moitié\n"
    "- France : leader sur la seconde main, 13 millions de membres, engagement environnemental\n"
    "- Innovation via IA et accélération digitale"
)

# 3. Business model de Rakuten
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "2. Business model de Rakuten"
content.text = (
    "- E-commerce : marketplace, neuf, occasion, reconditionné\n"
    "- Cashback et fidélité : Club R, cashback sur achats\n"
    "- Services financiers : banque en ligne, crédits, assurance\n"
    "- Télécommunications : Rakuten Mobile (Japon)\n"
    "- Logistique : Rakuten Fulfillment Network\n"
    "- Engagement environnemental : économie circulaire, seconde main"
)

# 4. Principaux concurrents
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "3. Principaux concurrents"
content.text = (
    "Mondiaux :\n"
    "- Amazon\n"
    "- Alibaba\n"
    "- eBay\n"
    "\n"
    "Français :\n"
    "- Cdiscount\n"
    "- Fnac Marketplace\n"
    "- Leboncoin\n"
    "\n"
    "Européens :\n"
    "- Zalando\n"
    "- Vinted\n"
    "- Bol.com"
)

# 5. Synthèse
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "4. Synthèse"
content.text = (
    "Rakuten combine e-commerce, services financiers et télécoms.\n"
    "Croissance et rentabilité opérationnelle en hausse.\n"
    "Leader sur la seconde main et l’économie circulaire en France.\n"
    "Innovation et engagement environnemental forts."
)

# Sauvegarde de la présentation
prs.save("rakuten_synthèse.pptx")
print("La présentation 'rakuten_synthèse.pptx' a été générée.")
