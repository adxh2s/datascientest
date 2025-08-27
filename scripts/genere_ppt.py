from pptx import Presentation
from pptx.util import Inches

# Création d'une présentation vide
prs = Presentation()

# 1. Diapositive de titre
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Nom du projet"
subtitle.text = "Présentation du projet"

# 2. Diapositive Business Model Canvas
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout "Titre et contenu"
title = slide.shapes.title
title.text = "Business Model Canvas"

# Option 1 : Texte explicatif (à remplacer par un tableau ou une image si besoin)
content = slide.placeholders[1]
content.text = ("Insérer ici le Business Model Canvas.\n"
                "Pour un vrai BMC, il est recommandé d'ajouter un tableau ou une image.\n"
                "Voir exemple de template : https://neoschronos.com/download/business-model-canvas/ppt/")

# Option 2 : Ajouter une image du BMC (décommente si tu as le fichier)
# img_path = "business_model_canvas.png"
# left = top = Inches(1)
# slide.shapes.add_picture(img_path, left, top, width=Inches(8), height=Inches(5))

# 3. Diapositive Persona 1
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Layout "Titre et contenu"
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Persona 1"
content.text = ("Nom : ...\n"
                "Âge : ...\n"
                "Métier : ...\n"
                "Besoins : ...\n"
                "Attentes : ...")

# 4. Diapositive Persona 2
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Persona 2"
content.text = ("Nom : ...\n"
                "Âge : ...\n"
                "Métier : ...\n"
                "Besoins : ...\n"
                "Attentes : ...")

# 5. Diapositive de fin
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Fin"
subtitle.text = "Merci pour votre attention"

# Sauvegarde de la présentation
prs.save("presentation_projet.pptx")
print("La présentation a été générée : presentation_projet.pptx")
